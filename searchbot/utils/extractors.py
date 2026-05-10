"""Extract text from different document file types."""

import os
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


def extract_text_from_pdf(file_path):
    """Read text from a PDF file."""
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"PDF extraction error: {e}")
    return text


def extract_text_from_docx(file_path):
    """Read text from a Word .docx file."""
    text = ""
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
    except Exception as e:
        print(f"DOCX extraction error: {e}")
    return text


def extract_text_from_xlsx(file_path):
    """Read text from an Excel .xlsx file."""
    text = ""
    try:
        wb = load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                text += row_text + "\n"
    except Exception as e:
        print(f"XLSX extraction error: {e}")
    return text


def extract_text_from_pptx(file_path):
    """Read text from a PowerPoint .pptx file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX extraction error: {e}")
    return text


def extract_text_from_txt(file_path):
    """Read text from a plain .txt file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    except Exception as e:
        print(f"TXT extraction error: {e}")
    return text


# Map of file extension to extractor function
EXTRACTORS = {
    'pdf': extract_text_from_pdf,
    'docx': extract_text_from_docx,
    'xlsx': extract_text_from_xlsx,
    'pptx': extract_text_from_pptx,
    'txt': extract_text_from_txt,
}

SUPPORTED_EXTENSIONS = list(EXTRACTORS.keys())


def get_file_extension(filename):
    """Return the lowercase file extension without the dot."""
    return os.path.splitext(filename)[1].lower().replace('.', '')


def extract_text(file_path, file_type):
    """Extract text from a file using the right extractor."""
    file_type = file_type.lower().replace('.', '')
    extractor = EXTRACTORS.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")
    return extractor(file_path)